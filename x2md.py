"""
x2md - Convert documents (PDF, DOCX, PPTX, TXT) to Markdown.

Default mode: MarkItDown only (fast, no API call).
Claude mode: runs multiple backends in parallel, merges via Claude.
Firecrawl mode: uploads file to temp hosting, parses via Firecrawl Document Parsing API.

Usage:
    python x2md.py                       # markitdown only (default)
    python x2md.py --claude              # multi-backend + Claude merge
    python x2md.py --firecrawl           # Firecrawl document parsing

As a module:
    from x2md import convert2md
    markdown_text = convert2md("/path/to/file.pdf")                    # markitdown only
    markdown_text = convert2md("/path/to/file.pdf", claude=True)       # multi-backend + merge
    markdown_text = convert2md("/path/to/file.pdf", firecrawl=True)    # Firecrawl parsing
"""

from datetime import datetime
import os
import sys
import time
import argparse
import re
import subprocess
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed

_running_as_cli = __name__ == '__main__'
if _running_as_cli:
    ts_time = datetime.now().strftime('%H:%M:%S')
    print(f"\n---------- {ts_time} starting {os.path.basename(__file__)}")
    start_time = time.time()


# ---------------------------------------------------------------------------
# File type detection
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    '.pdf', '.docx', '.pptx', '.txt', '.rtf', '.html', '.htm',
    '.csv', '.xlsx', '.xls', '.epub', '.md', '.json', '.xml',
}

PDF_EXTENSIONS = {'.pdf'}
DOCX_EXTENSIONS = {'.docx'}
PPTX_EXTENSIONS = {'.pptx'}
TEXT_EXTENSIONS = {'.txt', '.md', '.csv', '.json', '.xml', '.rtf'}


def detect_file_type(file_path: str) -> str:
    """Return the lowercase extension of the file."""
    return Path(file_path).suffix.lower()


# Matches data URIs (base64-encoded images, fonts, etc.) and long hex/base64 blobs
_BASE64_PATTERN = re.compile(
    r'(?:'
    r'!\[[^\]]*\]\(data:[^)]+\)'           # ![alt](data:image/...;base64,...)
    r'|'
    r'src=["\']data:[^"\']+["\']'           # src="data:image/...;base64,..."
    r'|'
    r'url\(data:[^)]+\)'                    # url(data:image/...;base64,...)
    r'|'
    r'\(data:[\w/;,+=\-]+\)'               # (data:image/png;base64,iVBOR...)
    r'|'
    r'data:[\w/]+;base64,[A-Za-z0-9+/=]{100,}'  # bare data:...;base64,... strings
    r')',
    re.DOTALL,
)


def _strip_base64(text: str) -> str:
    """Remove embedded base64 data URIs and long binary blobs from markdown."""
    cleaned = _BASE64_PATTERN.sub('[embedded image removed]', text)
    return cleaned


# ---------------------------------------------------------------------------
# Converter backends
# ---------------------------------------------------------------------------

def convert_markitdown(file_path: str) -> str | None:
    """Microsoft MarkItDown - broad format support, clean LLM-friendly output."""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(file_path)
        text = result.text_content
        if text and text.strip():
            return _strip_base64(text.strip())
    except Exception as e:
        print(f"  ❌ [markitdown] failed: {e}", file=sys.stderr)
    return None


def convert_pymupdf4llm(file_path: str) -> str | None:
    """PyMuPDF4LLM - LLM-optimized markdown from PDFs with layout awareness."""
    try:
        import pymupdf4llm
        text = pymupdf4llm.to_markdown(
            file_path,
            show_progress=False,
            force_text=True,
        )
        if text and text.strip():
            return _strip_base64(text.strip())
    except Exception as e:
        print(f"  ❌ [pymupdf4llm] failed: {e}", file=sys.stderr)
    return None


def convert_mammoth(file_path: str) -> str | None:
    """Mammoth - semantic DOCX-to-markdown, preserves document structure."""
    try:
        import mammoth
        with open(file_path, 'rb') as f:
            result = mammoth.convert_to_markdown(f)
            text = result.value
            if text and text.strip():
                return _strip_base64(text.strip())
    except Exception as e:
        print(f"  ❌ [mammoth] failed: {e}", file=sys.stderr)
    return None


def convert_pdfplumber(file_path: str) -> str | None:
    """pdfplumber - excellent table extraction from PDFs."""
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_parts = []

                # Extract tables
                tables = page.extract_tables()
                table_bboxes = []
                if tables:
                    for tbl in page.find_tables():
                        table_bboxes.append(tbl.bbox)

                    for table in tables:
                        if not table:
                            continue
                        # Build markdown table
                        cleaned = []
                        for row in table:
                            cleaned.append([
                                (cell or '').strip().replace('\n', ' ')
                                for cell in row
                            ])
                        if not cleaned:
                            continue
                        header = cleaned[0]
                        md_table = '| ' + ' | '.join(header) + ' |\n'
                        md_table += '| ' + ' | '.join(['---'] * len(header)) + ' |\n'
                        for row in cleaned[1:]:
                            # Pad row if shorter than header
                            while len(row) < len(header):
                                row.append('')
                            md_table += '| ' + ' | '.join(row[:len(header)]) + ' |\n'
                        page_parts.append(md_table)

                # Extract text (outside tables where possible)
                text = page.extract_text()
                if text and text.strip():
                    page_parts.insert(0, text.strip())

                if page_parts:
                    parts.append(f"<!-- Page {i + 1} -->\n\n" + '\n\n'.join(page_parts))

        if parts:
            return _strip_base64('\n\n---\n\n'.join(parts))
    except Exception as e:
        print(f"  ❌ [pdfplumber] failed: {e}", file=sys.stderr)
    return None


def convert_pptx_native(file_path: str) -> str | None:
    """python-pptx - direct slide-by-slide extraction with structure."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(file_path)
        parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"## Slide {slide_num}"]

            # Get slide title
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_parts[0] = f"## Slide {slide_num}: {slide.shapes.title.text.strip()}"

            # Extract notes
            notes_text = None
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        # Detect bullet level
                        level = para.level or 0
                        if level > 0:
                            indent = '  ' * (level - 1)
                            slide_parts.append(f"{indent}- {text}")
                        else:
                            # Skip if it's the title (already captured)
                            if shape == slide.shapes.title:
                                continue
                            slide_parts.append(text)

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        header = rows[0]
                        md_table = '| ' + ' | '.join(header) + ' |\n'
                        md_table += '| ' + ' | '.join(['---'] * len(header)) + ' |\n'
                        for row in rows[1:]:
                            while len(row) < len(header):
                                row.append('')
                            md_table += '| ' + ' | '.join(row[:len(header)]) + ' |\n'
                        slide_parts.append(md_table)

            if notes_text:
                slide_parts.append(f"\n> **Notes:** {notes_text}")

            parts.append('\n\n'.join(slide_parts))

        if parts:
            return _strip_base64('\n\n---\n\n'.join(parts))
    except Exception as e:
        print(f"  ❌ [pptx-native] failed: {e}", file=sys.stderr)
    return None


def convert_docx_native(file_path: str) -> str | None:
    """python-docx - direct paragraph/table extraction preserving styles."""
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        doc = Document(file_path)
        parts = []

        for element in doc.element.body:
            tag = element.tag.split('}')[-1]

            if tag == 'p':
                # It's a paragraph
                for para in [p for p in doc.paragraphs if p._element is element]:
                    text = para.text.strip()
                    if not text:
                        continue
                    style_name = (para.style.name or '').lower()

                    if 'heading 1' in style_name:
                        parts.append(f"# {text}")
                    elif 'heading 2' in style_name:
                        parts.append(f"## {text}")
                    elif 'heading 3' in style_name:
                        parts.append(f"### {text}")
                    elif 'heading 4' in style_name:
                        parts.append(f"#### {text}")
                    elif 'heading' in style_name:
                        parts.append(f"##### {text}")
                    elif 'list' in style_name or 'bullet' in style_name:
                        level = para.paragraph_format.left_indent
                        if level and level > 0:
                            parts.append(f"  - {text}")
                        else:
                            parts.append(f"- {text}")
                    else:
                        # Check for bold/italic runs
                        formatted = _format_docx_runs(para)
                        parts.append(formatted)

            elif tag == 'tbl':
                # It's a table
                for table in [t for t in doc.tables if t._element is element]:
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        header = rows[0]
                        md_table = '| ' + ' | '.join(header) + ' |\n'
                        md_table += '| ' + ' | '.join(['---'] * len(header)) + ' |\n'
                        for row in rows[1:]:
                            while len(row) < len(header):
                                row.append('')
                            md_table += '| ' + ' | '.join(row[:len(header)]) + ' |\n'
                        parts.append(md_table)

        if parts:
            return _strip_base64('\n\n'.join(parts))
    except Exception as e:
        print(f"  ❌ [docx-native] failed: {e}", file=sys.stderr)
    return None


def _format_docx_runs(para) -> str:
    """Format a docx paragraph's runs with bold/italic markdown."""
    parts = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        if run.bold and run.italic:
            text = f"***{text}***"
        elif run.bold:
            text = f"**{text}**"
        elif run.italic:
            text = f"*{text}*"
        parts.append(text)
    result = ''.join(parts).strip()
    return result if result else para.text.strip()


FIRECRAWL_SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.odt', '.rtf', '.xlsx', '.xls'}


def convert_firecrawl(file_path: str) -> str | None:
    """Firecrawl Document Parsing API - uploads to temp hosting, parses to markdown."""
    try:
        import requests
        from firecrawl import Firecrawl

        api_key = os.environ.get('FIRECRAWL_API_KEY_MARKDOWNEE')
        if not api_key:
            print("  ❌ [firecrawl] FIRECRAWL_API_KEY_MARKDOWNEE not set", file=sys.stderr)
            return None

        file_name = os.path.basename(file_path)
        print(f"  📤 [firecrawl] uploading {file_name} to temporary hosting...")
        with open(file_path, 'rb') as f:
            upload_resp = requests.post(
                'https://file.io',
                files={'file': (file_name, f)},
                timeout=120,
            )

        if not upload_resp.ok:
            print(f"  ❌ [firecrawl] file upload failed: {upload_resp.status_code}", file=sys.stderr)
            return None

        file_url = upload_resp.json().get('link')
        if not file_url:
            print(f"  ❌ [firecrawl] no URL returned from file upload", file=sys.stderr)
            return None

        print(f"  🔗 [firecrawl] uploaded, scraping via Firecrawl...")

        fc = Firecrawl(api_key=api_key)
        ext = Path(file_path).suffix.lower()
        parsers = None
        if ext == '.pdf':
            parsers = [{"type": "pdf", "mode": "auto"}]

        result = fc.scrape(file_url, formats=['markdown'], parsers=parsers)
        text = result.markdown
        if text and text.strip():
            return _strip_base64(text.strip())
        print(f"  ⚠️ [firecrawl] returned empty markdown", file=sys.stderr)

    except Exception as e:
        print(f"  ❌ [firecrawl] failed: {e}", file=sys.stderr)
    return None


def convert_text_direct(file_path: str) -> str | None:
    """Direct read for text-based files."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        if text and text.strip():
            return text.strip()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                text = f.read()
            if text and text.strip():
                return _strip_base64(text.strip())
        except Exception as e:
            print(f"  ❌ [text-direct] failed: {e}", file=sys.stderr)
    except Exception as e:
        print(f"  ❌ [text-direct] failed: {e}", file=sys.stderr)
    return None


# ---------------------------------------------------------------------------
# Converter registry - maps file types to applicable backends
# ---------------------------------------------------------------------------

CONVERTERS = {
    '.pdf': [
        ('markitdown', convert_markitdown),
        ('pymupdf4llm', convert_pymupdf4llm),
        ('pdfplumber', convert_pdfplumber),
    ],
    '.docx': [
        ('markitdown', convert_markitdown),
        ('mammoth', convert_mammoth),
        ('docx-native', convert_docx_native),
    ],
    '.pptx': [
        ('markitdown', convert_markitdown),
        ('pptx-native', convert_pptx_native),
    ],
    '.txt': [('text-direct', convert_text_direct)],
    '.md': [('text-direct', convert_text_direct)],
    '.csv': [('text-direct', convert_text_direct), ('markitdown', convert_markitdown)],
    '.json': [('text-direct', convert_text_direct)],
    '.xml': [('text-direct', convert_text_direct)],
    '.rtf': [('markitdown', convert_markitdown)],
    '.html': [('markitdown', convert_markitdown)],
    '.htm': [('markitdown', convert_markitdown)],
    '.xlsx': [('markitdown', convert_markitdown)],
    '.xls': [('markitdown', convert_markitdown)],
    '.epub': [('markitdown', convert_markitdown)],
}


# ---------------------------------------------------------------------------
# Claude merge skill
# ---------------------------------------------------------------------------

MERGE_SYSTEM_PROMPT = """You are a document conversion specialist. You receive multiple Markdown
conversions of the same source document, each produced by a different conversion tool.

Your job: produce the single best Markdown version by merging the strengths of each input.

Rules:
- Preserve ALL content from the original document. Do not drop sections, paragraphs, or data.
- Use the version with the best heading structure (proper # hierarchy).
- Use the version with the best table formatting (aligned pipes, complete cells).
- Use the version with the best list formatting (consistent bullets, proper nesting).
- Preserve bold, italic, and other inline formatting where any version captured it.
- Remove conversion artifacts: extra whitespace, repeated headers, garbled text.
- If one version has content another is missing, include it.
- Output clean Markdown only. No commentary, no explanations, no wrapper.
- Do not add any content that wasn't in the original document.
- Use standard hyphens for lists, never em-dashes."""

def merge_with_claude(
    results: dict[str, str],
    file_name: str,
    model: str = "claude-opus-4-6",
) -> str | None:
    """Send multiple converter outputs to Claude and get a merged best version."""
    try:
        import anthropic
        api_key = os.environ.get('ANTHROPIC_API_KEY')
        if not api_key:
            print("  ❌ [claude-merge] ANTHROPIC_API_KEY not set", file=sys.stderr)
            return None
        client = anthropic.Anthropic(api_key=api_key)

        # Build the prompt with each converter's output
        sections = []
        for name, text in results.items():
            # Truncate very long outputs to stay within context limits
            truncated = text[:80_000] if len(text) > 80_000 else text
            sections.append(
                f"=== OUTPUT FROM: {name.upper()} ===\n\n{truncated}\n\n"
                f"=== END {name.upper()} ==="
            )

        user_prompt = (
            f"Source file: {file_name}\n\n"
            f"Below are {len(results)} different Markdown conversions of the same document. "
            f"Merge them into the single best Markdown output.\n\n"
            + "\n\n".join(sections)
        )

        # Check total size - if too large, use extended thinking or trim
        total_chars = len(user_prompt)
        print(f"  🤖 [claude-merge] sending {total_chars:,} chars from {len(results)} converters")

        raw_response = client.messages.with_raw_response.create(
            model=model,
            max_tokens=16384,
            system=MERGE_SYSTEM_PROMPT,
            messages=[{"role": "user", "content": user_prompt}],
        )
        request_id = raw_response.http_response.headers.get('request-id', 'unknown')
        response = raw_response.parse()
        input_tokens = response.usage.input_tokens
        output_tokens = response.usage.output_tokens
        print(f"  🆔 [claude-merge] request_id: {request_id}")
        cost = (input_tokens * 5 + output_tokens * 25) / 1_000_000
        print(f"  📊 [claude-merge] tokens: {input_tokens:,} in / {output_tokens:,} out ({input_tokens + output_tokens:,} total) - ${cost:.4f}")

        merged = response.content[0].text
        if merged and merged.strip():
            return merged.strip()

    except Exception as e:
        print(f"  ❌ [claude-merge] failed: {e}", file=sys.stderr)
    return None


# ---------------------------------------------------------------------------
# Core conversion logic
# ---------------------------------------------------------------------------

def convert2md(file_path: str, claude: bool = False, firecrawl: bool = False) -> str:
    """Convert a document to Markdown.

    claude=False (default): MarkItDown only, fast, no API call.
    claude=True: runs all backends in parallel, merges via Claude.
    firecrawl=True: Firecrawl Document Parsing API.
    """
    file_path = os.path.abspath(file_path)
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = detect_file_type(file_path)
    file_name = os.path.basename(file_path)

    if firecrawl:
        if ext not in FIRECRAWL_SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Firecrawl does not support {ext}\n"
                f"Supported: {', '.join(sorted(FIRECRAWL_SUPPORTED_EXTENSIONS))}"
            )
        print(f"  📄 Converting {file_name} ({ext}) with Firecrawl...")
        result = convert_firecrawl(file_path)
        if result:
            print(f"  ✅ [firecrawl] OK ({len(result):,} chars)")
            return result
        _notify("x2md", f"Firecrawl conversion failed for {file_name}")
        print(f"\n❌ Firecrawl conversion failed for {file_name}", file=sys.stderr)
        sys.exit(1)

    if ext not in CONVERTERS:
        raise ValueError(
            f"Unsupported file type: {ext}\n"
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    if not claude:
        # Default mode: MarkItDown only
        print(f"  📄 Converting {file_name} ({ext}) with markitdown...")
        result = convert_markitdown(file_path)
        if result:
            print(f"  ✅ [markitdown] OK ({len(result):,} chars)")
            return result
        _notify("x2md", f"MarkItDown conversion failed for {file_name}")
        print(f"\n❌ MarkItDown conversion failed for {file_name}", file=sys.stderr)
        sys.exit(1)

    # Claude mode: all backends in parallel + merge
    converters = CONVERTERS[ext]
    print(f"  📄 Converting {file_name} ({ext}) with {len(converters)} backend(s)...")

    results: dict[str, str] = {}
    with ThreadPoolExecutor(max_workers=len(converters)) as executor:
        future_to_name = {
            executor.submit(fn, file_path): name
            for name, fn in converters
        }
        for future in as_completed(future_to_name):
            name = future_to_name[future]
            try:
                output = future.result()
                if output:
                    results[name] = output
                    print(f"  ✅ [{name}] OK ({len(output):,} chars)")
                else:
                    print(f"  ⚠️ [{name}] returned empty")
            except Exception as e:
                print(f"  ❌ [{name}] crashed: {e}", file=sys.stderr)

    if not results:
        _notify("x2md", f"All converters failed for {file_name}")
        print(f"\n❌ All converters failed for {file_name}", file=sys.stderr)
        sys.exit(1)

    print(f"  🔀 Merging {len(results)} outputs with Claude...")
    merged = merge_with_claude(results, file_name)
    if merged:
        print(f"  ✅ [claude-merge] OK ({len(merged):,} chars)")
        return merged

    _notify("x2md", "Claude merge failed. No output produced.")
    print(f"\n❌ Claude merge failed. Exiting without output.", file=sys.stderr)
    sys.exit(1)


def convert_and_save(file_path: str, output_path: str | None = None, claude: bool = False, firecrawl: bool = False) -> str:
    """Convert a document and save the Markdown to a file.

    Returns the output file path.
    """
    markdown = convert2md(file_path, claude=claude, firecrawl=firecrawl)

    if output_path is None:
        base, _ = os.path.splitext(file_path)
        output_path = f"{base}.md"

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(markdown)

    return output_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _notify(title: str, message: str) -> None:
    """Send a macOS notification."""
    safe_msg = message.replace('\\', '\\\\').replace('"', '\\"')
    safe_title = title.replace('\\', '\\\\').replace('"', '\\"')
    subprocess.run([
        'osascript', '-e',
        f'display notification "{safe_msg}" with title "{safe_title}"'
    ], check=False)


def _get_clipboard() -> str:
    """Read the macOS clipboard contents."""
    result = subprocess.run(['pbpaste'], capture_output=True, text=True, check=False)
    return result.stdout.strip()


if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        description='Convert documents (PDF, DOCX, PPTX, TXT, etc.) to Markdown'
    )
    parser.add_argument('-o', '--output', help='Output .md file path (default: same name as input)')
    parser.add_argument('--claude', action='store_true', help='Use multiple backends + Claude merge (default: markitdown only)')
    parser.add_argument('--firecrawl', action='store_true', help='Use Firecrawl Document Parsing API')
    parser.add_argument('--no-open', action='store_true', help='Do not open the output file in VS Code')
    args = parser.parse_args()

    # Read file path from clipboard
    input_path = _get_clipboard().strip("'\"")

    if not input_path:
        msg = "Clipboard is empty. Copy a file path first."
        print(f"\n❌ {msg}", file=sys.stderr)
        _notify("x2md", msg)
        sys.exit(1)

    if not os.path.isfile(input_path):
        msg = f"Not a valid file path: {input_path[:80]}"
        print(f"\n❌ {msg}", file=sys.stderr)
        _notify("x2md", msg)
        sys.exit(1)

    try:
        output_file = convert_and_save(
            input_path,
            output_path=args.output,
            claude=args.claude,
            firecrawl=args.firecrawl,
        )
        print(f"\n✅ Markdown saved to: {output_file}")

        if not args.no_open:
            os.system(f"open -a 'Visual Studio Code' '{output_file}'")

    except Exception as e:
        print(f"\n❌ Error: {e}", file=sys.stderr)
        sys.exit(1)

    print('\n-------------------------------')
    run_time = round((time.time() - start_time), 3)
    if run_time < 1:
        print(f'\n{os.path.basename(__file__)} finished in {round(run_time*1000)}ms at {datetime.now().strftime("%H:%M:%S")}.\n')
    elif run_time < 60:
        print(f'\n{os.path.basename(__file__)} finished in {round(run_time)}s at {datetime.now().strftime("%H:%M:%S")}.\n')
    elif run_time < 3600:
        print(f'\n{os.path.basename(__file__)} finished in {round(run_time/60)}mns at {datetime.now().strftime("%H:%M:%S")}.\n')
    else:
        print(f'\n{os.path.basename(__file__)} finished in {round(run_time/3600, 2)}hrs at {datetime.now().strftime("%H:%M:%S")}.\n')
